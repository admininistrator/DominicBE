"""Knowledge ingestion service – chunking, indexing skeleton for RAG."""
import hashlib
import logging
import re
import time
from typing import Optional

from sqlalchemy.orm import Session

from app.core.config import settings
from app.crud import crud_knowledge

logger = logging.getLogger("uvicorn.error")
LOCAL_EMBEDDING_DIMENSIONS = 64


# ---------------------------------------------------------------------------
# Lazy image captioning helper (avoids circular import with llm_provider)
# ---------------------------------------------------------------------------

def _caption_image(image_bytes: bytes, media_type: str = "image/jpeg", context_hint: str = "") -> str:
    """Caption an image via llm_provider. Returns '' when captioning is disabled or on error."""
    if not settings.llm_image_captioning_enabled:
        return ""
    if not image_bytes:
        return ""
    try:
        from app.services.llm_provider import caption_image
        return caption_image(image_bytes, media_type=media_type, context_hint=context_hint)
    except Exception as exc:
        logger.warning("Image captioning failed (non-fatal): %s", exc)
        return ""


def normalize_text_for_ingestion(text: str) -> str:
    """Normalize whitespace while preserving readable paragraph boundaries."""
    raw = (text or "").replace("\r\n", "\n").replace("\r", "\n").strip()
    if not raw:
        return ""

    paragraphs = []
    for block in re.split(r"\n\s*\n+", raw):
        normalized = re.sub(r"[ \t]+", " ", block).strip()
        if normalized:
            paragraphs.append(normalized)
    return "\n\n".join(paragraphs)


def resolve_embedding_model_name() -> str:
    configured = (settings.embedding_model or "").strip()
    return configured or "local-hash-v1"


def compute_text_embedding(text: str, *, dimensions: int = LOCAL_EMBEDDING_DIMENSIONS) -> list[float]:
    normalized = normalize_text_for_ingestion(text).lower()
    if not normalized:
        return [0.0] * dimensions

    vector = [0.0] * dimensions
    for token in re.findall(r"\w+", normalized, flags=re.UNICODE):
        token_hash = hashlib.sha256(token.encode("utf-8")).digest()
        bucket = int.from_bytes(token_hash[:2], "big") % dimensions
        sign = 1.0 if token_hash[2] % 2 == 0 else -1.0
        weight = 1.0 + ((token_hash[3] % 5) * 0.1)
        vector[bucket] += sign * weight

    magnitude = sum(value * value for value in vector) ** 0.5
    if magnitude == 0:
        return [0.0] * dimensions
    return [round(value / magnitude, 6) for value in vector]


def build_vector_id(document_id: int, chunk_index: int, checksum: str) -> str:
    return f"local:{document_id}:{chunk_index}:{checksum[:12]}"


def prepare_chunks_for_indexing(document_id: int, checksum: str, chunks: list[dict]) -> list[dict]:
    embedding_model = resolve_embedding_model_name()
    prepared: list[dict] = []

    for chunk in chunks:
        content = chunk["content"]
        metadata_json = {
            **(chunk.get("metadata_json") or {}),
            "embedding": compute_text_embedding(content),
            "index_provider": settings.vector_store_provider,
            "embedding_provider": settings.embedding_provider,
        }
        prepared.append(
            {
                **chunk,
                "embedding_model": embedding_model,
                "vector_id": build_vector_id(document_id, chunk["chunk_index"], checksum),
                "metadata_json": metadata_json,
            }
        )

    return prepared


# ---------------------------------------------------------------------------
# Text chunking (real implementation)
# ---------------------------------------------------------------------------

def chunk_text(
    text: str,
    chunk_size: int | None = None,
    chunk_overlap: int | None = None,
) -> list[dict]:
    """Split text into overlapping chunks.

    Returns list of {chunk_index, content, token_count (estimated)}.
    Uses sentence-boundary-aware splitting when possible.
    """
    size = chunk_size or settings.chunk_size
    overlap = chunk_overlap or settings.chunk_overlap

    normalized_text = normalize_text_for_ingestion(text)
    if not normalized_text:
        return []

    # Split into sentences first for cleaner boundaries
    sentences = _split_sentences(normalized_text)
    chunks: list[dict] = []
    current_chunk: list[str] = []
    current_len = 0
    idx = 0

    expanded_sentences: list[str] = []
    for sentence in sentences:
        expanded_sentences.extend(_split_large_sentence(sentence, size=size, overlap=overlap))

    for sentence in expanded_sentences:
        s_len = len(sentence)
        if current_len + s_len > size and current_chunk:
            chunk_text_str = " ".join(current_chunk).strip()
            if chunk_text_str:
                chunks.append({
                    "chunk_index": idx,
                    "content": chunk_text_str,
                    "token_count": max(1, len(chunk_text_str) // 4),
                    "metadata_json": {"char_count": len(chunk_text_str)},
                })
                idx += 1

            # Keep overlap: walk back from end
            overlap_chunks: list[str] = []
            overlap_len = 0
            for s in reversed(current_chunk):
                if overlap_len + len(s) > overlap:
                    break
                overlap_chunks.insert(0, s)
                overlap_len += len(s)
            current_chunk = overlap_chunks
            current_len = overlap_len

        current_chunk.append(sentence)
        current_len += s_len

    # Last chunk
    if current_chunk:
        chunk_text_str = " ".join(current_chunk).strip()
        if chunk_text_str:
            chunks.append({
                "chunk_index": idx,
                "content": chunk_text_str,
                "token_count": max(1, len(chunk_text_str) // 4),
                "metadata_json": {"char_count": len(chunk_text_str)},
            })

    return chunks


def _split_sentences(text: str) -> list[str]:
    """Naive sentence splitter: split on '. ', '! ', '? ', newlines."""
    # Split on sentence-ending punctuation followed by space or newline
    parts = re.split(r'(?<=[.!?])\s+|\n+', text)
    return [p.strip() for p in parts if p.strip()]


def _split_large_sentence(sentence: str, size: int, overlap: int) -> list[str]:
    if len(sentence) <= size:
        return [sentence]

    parts: list[str] = []
    start = 0
    step = max(1, size - overlap)
    while start < len(sentence):
        end = min(len(sentence), start + size)
        piece = sentence[start:end].strip()
        if piece:
            parts.append(piece)
        if end >= len(sentence):
            break
        start += step
    return parts or [sentence]


# ---------------------------------------------------------------------------
# File text extraction
# ---------------------------------------------------------------------------

def extract_text_from_file(content: bytes, filename: str, mime_type: str | None = None) -> str:
    """Extract plain text from uploaded file.

    Supported formats:
    - txt, md, csv, log, json, py, js, yaml/yml : decode as UTF-8
    - pdf  : PyMuPDF  – text blocks + embedded table heuristics
    - docx : python-docx – paragraphs + tables + text-boxes
    - pptx : python-pptx – all slide shapes, speaker notes, tables
    - xlsx : openpyxl   – all sheets, all non-empty cells (header|value rows)
    """
    lower = filename.lower()

    # ── Plain text formats ────────────────────────────────────────────────────
    if lower.endswith((".txt", ".md", ".csv", ".log", ".json", ".py", ".js", ".yaml", ".yml")):
        return content.decode("utf-8", errors="replace")

    # ── PDF ──────────────────────────────────────────────────────────────────
    if lower.endswith(".pdf"):
        return _extract_pdf(content)

    # ── DOCX ─────────────────────────────────────────────────────────────────
    if lower.endswith(".docx"):
        return _extract_docx(content)

    # ── PPTX ─────────────────────────────────────────────────────────────────
    if lower.endswith(".pptx"):
        return _extract_pptx(content)

    # ── XLSX / XLS ───────────────────────────────────────────────────────────
    if lower.endswith((".xlsx", ".xls")):
        return _extract_xlsx(content)

    raise ValueError(f"Unsupported file type: {filename}")


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF using PyMuPDF.

    Strategy:
    1. Use get_text("blocks") to preserve reading order (sorted by y, x).
    2. When LLM_IMAGE_CAPTIONING_ENABLED=true, extract embedded images per page
       and inject AI-generated captions as [Image N: <description>] markers.
    3. Page numbers are prepended as markers so chunk metadata can reference them.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ValueError("PDF support requires PyMuPDF. Install with: pip install PyMuPDF")

    doc = fitz.open(stream=content, filetype="pdf")
    page_texts: list[str] = []

    for page_num, page in enumerate(doc, start=1):
        blocks = page.get_text("blocks")  # (x0,y0,x1,y1, text, block_no, block_type)
        blocks_sorted = sorted(
            [b for b in blocks if b[6] == 0 and b[4].strip()],  # type==0 → text block
            key=lambda b: (round(b[1] / 10) * 10, b[0]),
        )
        parts = [f"[Page {page_num}]"] + [b[4].strip() for b in blocks_sorted]

        # Extract and caption embedded images when enabled
        if settings.llm_image_captioning_enabled:
            for img_idx, img_info in enumerate(page.get_images(full=True), start=1):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    ext = (base_image.get("ext") or "jpeg").lower().replace("jpg", "jpeg")
                    mt = f"image/{ext}" if ext in ("jpeg", "png", "gif", "webp") else "image/jpeg"
                    caption = _caption_image(img_bytes, media_type=mt, context_hint=f"PDF page {page_num}")
                    if caption:
                        parts.append(f"[Image {img_idx}: {caption}]")
                except Exception as exc:
                    logger.debug("PDF image extraction skipped xref=%s: %s", xref, exc)

        page_texts.append("\n".join(parts))

    doc.close()
    return "\n\n".join(page_texts)


def _extract_docx(content: bytes) -> str:
    """Extract text from DOCX using python-docx with full XML traversal."""
    try:
        import io
        from docx import Document
    except ImportError:
        raise ValueError("DOCX support requires python-docx. Install with: pip install python-docx")

    WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    def _para_full_text(para) -> str:
        texts = []
        for node in para._element.iter(f"{{{WORD_NS}}}t"):
            t = node.text or ""
            if t:
                texts.append(t)
        return "".join(texts)

    def _extract_table(table) -> list[str]:
        rows: list[str] = []
        for row in table.rows:
            row_parts: list[str] = []
            seen: set[str] = set()
            for cell in row.cells:
                cell_texts = [_para_full_text(p).strip() for p in cell.paragraphs]
                cell_text = "\n".join(t for t in cell_texts if t).strip()
                if cell_text and cell_text not in seen:
                    row_parts.append(cell_text)
                    seen.add(cell_text)
                for nested in cell.tables:
                    rows.extend(_extract_table(nested))
            if row_parts:
                rows.append(" | ".join(row_parts))
        return rows

    doc = Document(io.BytesIO(content))
    parts: list[str] = []

    # Paragraphs
    for para in doc.paragraphs:
        text = _para_full_text(para).strip()
        if text:
            parts.append(text)

    # Tables
    for table in doc.tables:
        parts.extend(_extract_table(table))

    # Embedded images – caption when enabled
    if settings.llm_image_captioning_enabled:
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_bytes = rel.target_part.blob
                    mt = rel.target_part.content_type or "image/jpeg"
                    caption = _caption_image(img_bytes, media_type=mt, context_hint="DOCX document")
                    if caption:
                        parts.append(f"[Image: {caption}]")
                except Exception as exc:
                    logger.debug("DOCX image extraction skipped: %s", exc)

    return "\n".join(parts)


def _extract_pptx(content: bytes) -> str:
    """Extract text from PPTX using python-pptx.

    Extracts per slide:
    - Title (prefixed with slide number and title)
    - All text frames (body text, text boxes, grouped shapes)
    - Tables (formatted as header | value rows)
    - Speaker notes
    """
    try:
        import io
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise ValueError("PPTX support requires python-pptx. Install with: pip install python-pptx")

    prs = Presentation(io.BytesIO(content))
    slide_texts: list[str] = []

    def _shape_texts(shape) -> list[str]:
        """Recursively extract text from a shape (including groups)."""
        results: list[str] = []
        # Group shape: recurse
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                results.extend(_shape_texts(child))
            return results
        # Table
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_cells:
                    results.append(" | ".join(row_cells))
            return results
        # Text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    results.append(text)
        return results

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        # Slide title
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text.strip()

        header = f"[Slide {slide_num}]"
        if title_text:
            header += f" {title_text}"
        parts.append(header)

        # All shapes
        for shape in slide.shapes:
            # Skip re-reading the title shape to avoid duplication
            if shape == slide.shapes.title:
                continue
            parts.extend(_shape_texts(shape))

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        # Embedded images on slide – caption when enabled
        if settings.llm_image_captioning_enabled:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        mt = shape.image.content_type or "image/jpeg"
                        caption = _caption_image(img_bytes, media_type=mt, context_hint=f"Slide {slide_num}: {title_text}")
                        if caption:
                            parts.append(f"[Image: {caption}]")
                    except Exception as e:
                        logger.debug("PPTX image extraction skipped slide=%s: %s", slide_num, e)

        slide_texts.append("\n".join(parts))

    return "\n\n".join(slide_texts)


def _extract_xlsx(content: bytes) -> str:
    """Extract text from XLSX/XLS using openpyxl.

    Strategy:
    - For each worksheet: emit sheet name as section header.
    - Detect header row (first non-empty row) and format subsequent rows as
      "Header1: value1 | Header2: value2 …" for semantic richness.
    - Falls back to raw cell values when headers are numeric/ambiguous.
    - Merged cells are read from their master cell value.
    """
    try:
        import io
        import openpyxl
    except ImportError:
        raise ValueError("XLSX support requires openpyxl. Install with: pip install openpyxl")

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    sheet_texts: list[str] = []

    for sheet in wb.worksheets:
        parts: list[str] = [f"[Sheet: {sheet.title}]"]
        rows_data: list[list[str]] = []

        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() if cell is not None else "" for cell in row]
            if any(cells):  # skip fully empty rows
                rows_data.append(cells)

        if not rows_data:
            continue

        # Use first row as headers if it looks like text
        first_row = rows_data[0]
        is_header_row = any(cell and not _is_numeric(cell) for cell in first_row)

        if is_header_row and len(rows_data) > 1:
            headers = first_row
            for row_cells in rows_data[1:]:
                pairs = []
                for h, v in zip(headers, row_cells):
                    if v:
                        label = h if h and not _is_numeric(h) else ""
                        pairs.append(f"{label}: {v}" if label else v)
                if pairs:
                    parts.append(" | ".join(pairs))
        else:
            # No clear header row – just join cells
            for row_cells in rows_data:
                line = " | ".join(c for c in row_cells if c)
                if line:
                    parts.append(line)

        sheet_texts.append("\n".join(parts))

    wb.close()
    return "\n\n".join(sheet_texts)


def _is_numeric(value: str) -> bool:
    try:
        float(value)
        return True
    except ValueError:
        return False


# ---------------------------------------------------------------------------
# Shared internal indexing logic (avoids duplication across 3 callers)
# ---------------------------------------------------------------------------

def _execute_indexing(db: Session, doc_id: int, job_id: int) -> dict:
    """Core chunk → embed → store logic. Updates job/document status.

    Raises on failure (caller is responsible for error-state updates).
    Returns result dict on success.
    """
    doc = crud_knowledge.get_document(db, doc_id)
    if not doc:
        raise ValueError(f"Document {doc_id} not found.")
    if not doc.raw_text:
        raise ValueError(f"Document {doc_id} has no raw text.")

    normalized_text = normalize_text_for_ingestion(doc.raw_text)
    if not normalized_text:
        raise ValueError(f"Document {doc_id} has no valid text after normalization.")

    checksum = hashlib.sha256(normalized_text.encode("utf-8")).hexdigest()

    crud_knowledge.update_document_status(db, doc_id, "processing")
    crud_knowledge.update_ingestion_job_status(db, job_id, "processing")

    chunks = chunk_text(normalized_text)
    if not chunks:
        raise ValueError("Document produced no chunks after splitting.")

    logger.info(
        "Indexing doc=%s job=%s: %d chunks (avg %d chars)",
        doc_id, job_id, len(chunks),
        sum(len(c["content"]) for c in chunks) // max(len(chunks), 1),
    )

    indexed_chunks = prepare_chunks_for_indexing(doc_id, checksum, chunks)
    crud_knowledge.delete_chunks_by_document(db, doc_id)
    chunk_rows = crud_knowledge.create_chunks_bulk(db, doc_id, indexed_chunks)

    logger.info(
        "Indexed doc=%s job=%s embedding_provider=%s model=%s vector_store=%s",
        doc_id, job_id,
        settings.embedding_provider, resolve_embedding_model_name(), settings.vector_store_provider,
    )

    crud_knowledge.update_document_status(db, doc_id, "indexed")
    crud_knowledge.update_ingestion_job_status(db, job_id, "completed")

    return {
        "document_id": doc_id,
        "job_id": job_id,
        "status": "indexed",
        "chunks_count": len(chunk_rows),
        "checksum": checksum,
    }


# ---------------------------------------------------------------------------
# Ingestion pipeline (synchronous for now, async/background later)
# ---------------------------------------------------------------------------

def create_document_record(
    db: Session,
    owner_username: str,
    title: str,
    raw_text: str,
    source_type: str = "text",
    source_uri: Optional[str] = None,
    mime_type: Optional[str] = None,
    metadata: Optional[dict] = None,
) -> dict:
    """Create document + job records synchronously and return IDs immediately.

    Call ``run_indexing_pipeline`` afterwards (e.g. as a background task) to
    perform the actual chunking / embedding / indexing work.
    """
    normalized_owner = (owner_username or "").strip()
    normalized_title = (title or "").strip()
    normalized_text = normalize_text_for_ingestion(raw_text)

    if not normalized_owner:
        raise ValueError("owner_username is required.")
    if not normalized_title:
        raise ValueError("title is required.")
    if not normalized_text:
        raise ValueError("raw_text is empty after normalization.")

    checksum = hashlib.sha256(normalized_text.encode("utf-8")).hexdigest()

    existing = crud_knowledge.get_document_by_owner_and_checksum(db, normalized_owner, checksum)
    if existing:
        logger.info(
            "Document id=%s already exists for owner=%s checksum=%s – creating new job for reindex.",
            existing.id, normalized_owner, checksum[:12],
        )
        job = crud_knowledge.create_ingestion_job(db, existing.id)
        return {"document_id": existing.id, "job_id": job.id, "status": "pending", "checksum": checksum}

    doc = crud_knowledge.create_document(
        db=db,
        owner_username=normalized_owner,
        title=normalized_title,
        source_type=source_type,
        source_uri=source_uri,
        mime_type=mime_type,
        raw_text=normalized_text,
        checksum=checksum,
        metadata_json=metadata,
    )
    job = crud_knowledge.create_ingestion_job(db, doc.id)
    return {"document_id": doc.id, "job_id": job.id, "status": "pending", "checksum": checksum}


def run_indexing_pipeline(doc_id: int, job_id: int, db_factory) -> dict:
    """Background-safe wrapper around ``_execute_indexing`` with retry.

    ``db_factory`` must be a callable that returns a fresh ``Session``.
    Configured via ``INGESTION_MAX_RETRIES`` / ``INGESTION_RETRY_DELAY_SECONDS``.
    """
    max_attempts = max(1, settings.ingestion_max_retries + 1)
    delay = settings.ingestion_retry_delay_seconds
    last_exc: Exception | None = None

    for attempt in range(1, max_attempts + 1):
        db = db_factory()
        try:
            result = _execute_indexing(db, doc_id, job_id)
            result["attempts"] = attempt
            return result
        except Exception as e:
            last_exc = e
            logger.warning(
                "Indexing attempt %s/%s failed doc=%s job=%s: %s",
                attempt, max_attempts, doc_id, job_id, e,
            )
            try:
                crud_knowledge.update_document_status(db, doc_id, "failed")
                crud_knowledge.update_ingestion_job_status(db, job_id, "failed", str(e))
            except Exception:
                pass

            if attempt < max_attempts:
                backoff = delay * (2 ** (attempt - 1))
                logger.info("Retrying doc=%s in %.1fs ...", doc_id, backoff)
                time.sleep(backoff)
        finally:
            db.close()

    logger.error(
        "All %s indexing attempts exhausted for doc=%s job=%s. Last error: %s",
        max_attempts, doc_id, job_id, last_exc,
    )
    raise last_exc


def ingest_document(
    db: Session,
    owner_username: str,
    title: str,
    raw_text: str,
    source_type: str = "text",
    source_uri: Optional[str] = None,
    mime_type: Optional[str] = None,
    metadata: Optional[dict] = None,
) -> dict:
    """Full synchronous ingestion: create doc → chunk → index.

    Use ``create_document_record`` + ``run_indexing_pipeline`` for async flows.
    """
    normalized_owner = (owner_username or "").strip()
    normalized_title = (title or "").strip()
    normalized_text = normalize_text_for_ingestion(raw_text)

    if not normalized_owner:
        raise ValueError("owner_username is required.")
    if not normalized_title:
        raise ValueError("title is required.")
    if not normalized_text:
        raise ValueError("raw_text is empty after normalization.")

    checksum = hashlib.sha256(normalized_text.encode("utf-8")).hexdigest()

    existing = crud_knowledge.get_document_by_owner_and_checksum(db, normalized_owner, checksum)
    if existing:
        logger.info(
            "Reusing existing knowledge document id=%s for owner=%s checksum=%s",
            existing.id, normalized_owner, checksum[:12],
        )
        return reindex_document(db, existing.id)

    doc = crud_knowledge.create_document(
        db=db,
        owner_username=normalized_owner,
        title=normalized_title,
        source_type=source_type,
        source_uri=source_uri,
        mime_type=mime_type,
        raw_text=normalized_text,
        checksum=checksum,
        metadata_json=metadata,
    )
    job = crud_knowledge.create_ingestion_job(db, doc.id)

    try:
        return _execute_indexing(db, doc.id, job.id)
    except Exception as e:
        logger.error("Ingestion failed for doc %d: %s", doc.id, e)
        raise


def ingest_uploaded_file(
    db: Session,
    owner_username: str,
    filename: str,
    content: bytes,
    mime_type: str | None = None,
) -> dict:
    """Convenience: extract text from file bytes, then run full ingestion."""
    raw_text = extract_text_from_file(content, filename, mime_type)
    return ingest_document(
        db=db,
        owner_username=owner_username,
        title=filename,
        raw_text=raw_text,
        source_type="upload",
        source_uri=filename,
        mime_type=mime_type,
    )


def reindex_document(db: Session, doc_id: int) -> dict:
    """Re-chunk and re-index an existing document."""
    doc = crud_knowledge.get_document(db, doc_id)
    if not doc:
        raise ValueError("Document not found.")
    if not doc.raw_text:
        raise ValueError("Document has no raw text to re-index.")

    job = crud_knowledge.create_ingestion_job(db, doc.id)
    try:
        return _execute_indexing(db, doc_id, job.id)
    except Exception as e:
        logger.error("Reindex failed for doc %d: %s", doc_id, e)
        raise

